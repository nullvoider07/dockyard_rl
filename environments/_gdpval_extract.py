# Standalone deliverable extractor for the GDPval agentic environment.
#
# Runs INSIDE the sandbox container (no dockyard imports — pure stdlib + optional
# document libraries). The agentic environment ships this file's source into the
# container at grading time and runs `python3 <this> <dir> [char_limit]`; its
# stdout becomes the deliverable text graded by GDPvalRubricReward.
#
# It walks a directory, emits a file manifest, then a best-effort text extraction
# per file type. Every per-file extractor is guarded so a missing library or a
# corrupt file degrades to a placeholder line rather than aborting the run.

import os
import sys

DEFAULT_CHAR_LIMIT = 200_000

# Extensions read directly as UTF-8 text.
TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst", ".csv", ".tsv", ".json", ".jsonl",
    ".html", ".htm", ".xml", ".yaml", ".yml", ".py", ".tex", ".log", ".ini",
    ".cfg", ".toml", ".sql", ".r", ".m", ".sh",
}


def _read_text(path):
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as fh:
            return fh.read()
    except Exception as exc:  # noqa: BLE001
        return f"[could not read text: {exc}]"


def _extract_xlsx(path):
    try:
        from openpyxl import load_workbook
    except Exception as exc:  # noqa: BLE001
        return f"[xlsx: openpyxl unavailable: {exc}]"
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f"[sheet: {ws.title}]")
            for row in ws.iter_rows(values_only=True):
                cells = ["" if c is None else str(c) for c in row]
                if any(cells):
                    parts.append("\t".join(cells))
        wb.close()
        return "\n".join(parts)
    except Exception as exc:  # noqa: BLE001
        return f"[could not extract xlsx: {exc}]"


def _extract_docx(path):
    try:
        import docx  # python-docx
    except Exception as exc:  # noqa: BLE001
        return f"[docx: python-docx unavailable: {exc}]"
    try:
        document = docx.Document(path)
        parts = [p.text for p in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                parts.append("\t".join(cell.text for cell in row.cells))
        return "\n".join(parts)
    except Exception as exc:  # noqa: BLE001
        return f"[could not extract docx: {exc}]"


def _extract_pptx(path):
    try:
        from pptx import Presentation
    except Exception as exc:  # noqa: BLE001
        return f"[pptx: python-pptx unavailable: {exc}]"
    try:
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"[slide {i}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs)
                        if text:
                            parts.append(text)
        return "\n".join(parts)
    except Exception as exc:  # noqa: BLE001
        return f"[could not extract pptx: {exc}]"


def _extract_pdf(path):
    try:
        from pypdf import PdfReader
    except Exception as exc:  # noqa: BLE001
        return f"[pdf: pypdf unavailable: {exc}]"
    try:
        reader = PdfReader(path)
        return "\n".join((page.extract_text() or "") for page in reader.pages)
    except Exception as exc:  # noqa: BLE001
        return f"[could not extract pdf: {exc}]"


_BINARY_EXTRACTORS = {
    ".xlsx": _extract_xlsx,
    ".xlsm": _extract_xlsx,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".pdf": _extract_pdf,
}


def extract_file(path):
    """Best-effort text extraction for a single file, dispatched by extension."""
    ext = os.path.splitext(path)[1].lower()
    if ext in TEXT_EXTENSIONS:
        return _read_text(path)
    extractor = _BINARY_EXTRACTORS.get(ext)
    if extractor is not None:
        return extractor(path)
    return f"[binary file, no text extractor for '{ext}']"


def walk_files(root):
    """Return (relpath, abspath, size) for every regular file under root, sorted."""
    found = []
    for dirpath, _dirs, names in os.walk(root):
        for name in sorted(names):
            ap = os.path.join(dirpath, name)
            if not os.path.isfile(ap):
                continue
            rel = os.path.relpath(ap, root)
            try:
                size = os.path.getsize(ap)
            except OSError:
                size = -1
            found.append((rel, ap, size))
    found.sort(key=lambda t: t[0])
    return found


def extract_dir(root, char_limit=DEFAULT_CHAR_LIMIT):
    """Build the manifest + extracted-content bundle for a deliverable directory."""
    if not os.path.isdir(root):
        return f"[no deliverable directory at {root}]"
    files = walk_files(root)
    if not files:
        return f"[deliverable directory {root} is empty — no files produced]"

    lines = ["=== DELIVERABLE MANIFEST ==="]
    for rel, _ap, size in files:
        ext = os.path.splitext(rel)[1].lower() or "(none)"
        lines.append(f"- {rel}  ({size} bytes, type {ext})")
    lines.append("")
    lines.append("=== EXTRACTED CONTENT ===")

    out = "\n".join(lines)
    for rel, ap, _size in files:
        if len(out) >= char_limit:
            out += "\n... [truncated: remaining files omitted]"
            break
        body = extract_file(ap)
        section = f"\n\n--- FILE: {rel} ---\n{body}"
        remaining = char_limit - len(out)
        if len(section) > remaining:
            section = section[:remaining] + "\n... [truncated]"
        out += section
    return out


def main(argv):
    if len(argv) < 2:
        print("[extractor error: no directory argument]")
        return 0
    root = argv[1]
    char_limit = DEFAULT_CHAR_LIMIT
    if len(argv) >= 3:
        try:
            char_limit = int(argv[2])
        except (TypeError, ValueError):
            char_limit = DEFAULT_CHAR_LIMIT
    sys.stdout.write(extract_dir(root, char_limit))
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
